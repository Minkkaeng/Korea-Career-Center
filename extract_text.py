import os
import glob
import docx
import pptx
import fitz

artifact_dir = "C:/Users/choi9/.gemini/antigravity/brain/41d84f6a-2431-4faf-aa73-84bc7b9a4072"
base_dir = "./전달자료"

def extract_docx(path):
    doc = docx.Document(path)
    return "\n".join([p.text for p in doc.paragraphs])

def extract_pptx(path):
    prs = pptx.Presentation(path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def extract_pdf(path):
    doc = fitz.open(path)
    text = []
    for page in doc:
        text.append(page.get_text())
    return "\n".join(text)

files = glob.glob(f"{base_dir}/**/*.*", recursive=True)
summary = []

for f in files:
    ext = f.lower().split('.')[-1]
    content = ""
    try:
        if ext == 'docx':
            content = extract_docx(f)
        elif ext == 'pptx':
            content = extract_pptx(f)
        elif ext == 'pdf':
            content = extract_pdf(f)
        elif ext == 'txt':
            with open(f, 'r', encoding='utf-8') as txt_f:
                content = txt_f.read()
        
        if content:
            name = os.path.basename(f)
            out_path = os.path.join(artifact_dir, f"{name}.md")
            with open(out_path, 'w', encoding='utf-8') as out_f:
                out_f.write(content)
            summary.append(f"Extracted: {name}")
    except Exception as e:
        summary.append(f"Error on {f}: {e}")

print("\n".join(summary))
